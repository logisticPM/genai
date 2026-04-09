"""FinChartAudit poster v16 — experimental methodology, headline finding,
conclusory captions, <400 words, clean template."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation('C:/Users/chntw/Documents/7180/DD_v1/Capstone Template 4-1.pptx')
slide = prs.slides[0]

BLUE  = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED   = RGBColor(0xC6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x33, 0x33, 0x33)
GRAY  = RGBColor(0x66, 0x66, 0x66)

SZ_TITLE=50; SZ_HEAD=26; SZ_BODY=20; SZ_CAP=16; SZ_SUB=18; SZ_HERO=24

shapes={s.name:s for s in slide.shapes}
F='C:/Users/chntw/Documents/7180/DD_v1/results/figures'
IMG='C:/Users/chntw/Documents/7180/DD_v1/data/misviz/img'

def tb(l,t,w,h):
    b=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    b.text_frame.word_wrap=True; return b.text_frame
def wr(tf,txt,sz=SZ_BODY,bold=False,color=DARK,align=PP_ALIGN.LEFT,italic=False,sa=Pt(4)):
    p=tf.paragraphs[0] if(tf.paragraphs and tf.paragraphs[0].text=='' and len(tf.paragraphs)==1) else tf.add_paragraph()
    p.alignment,p.space_after,p.space_before=align,sa,Pt(0)
    r=p.add_run(); r.text=txt; r.font.size,r.font.bold,r.font.italic=Pt(sz),bold,italic
    r.font.color.rgb,r.font.name=color,'Calibri'
def bl(tf,txt,sz=SZ_BODY,color=DARK): wr(tf,'\u2022 '+txt,sz=sz,color=color,sa=Pt(4))
def hd(tf,txt,sz=SZ_HEAD): wr(tf,txt,sz=sz,bold=True,color=BLUE,sa=Pt(6))
def cap(tf,txt): wr(tf,txt,sz=SZ_CAP,italic=True,color=GRAY,sa=Pt(2))
def hide(n):
    s=shapes.get(n)
    if s: s.left=Inches(30)

for n in ['TextBox 26','TextBox 291','TextBox 1','Rectangle 8','Rectangle 116','Rectangle 117',
          'Rectangle 47','Rectangle 76','Rectangle 77','Rectangle 69','Rectangle 71','TextBox 288']:
    s=shapes.get(n)
    if s and s.has_text_frame: s.text_frame.clear()
hide('Rectangle 47'); hide('Rectangle 76'); hide('Rectangle 77')
hide('Freeform 19'); hide('Picture 292'); hide('Picture 58'); hide('Picture 329')

# ================================================================
# TITLE + AUTHORS
# ================================================================
shapes['TextBox 1'].text_frame.clear()
wr(shapes['TextBox 1'].text_frame,
   'FinChartAudit: Detecting Misleading Financial Charts with Vision-Language Models',
   sz=SZ_TITLE,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

tf_a=tb(2.50,1.05,18.50,0.50)
wr(tf_a,'Mengshan Li,  Wenshuang Zhou,  Tong Wu',sz=24,bold=True,color=DARK,align=PP_ALIGN.CENTER)
tf_af=tb(2.50,1.55,18.50,0.35)
wr(tf_af,'Khoury College of Computer Sciences, Northeastern University',sz=18,color=GRAY,align=PP_ALIGN.CENTER)

# ================================================================
# LEFT COLUMN
# ================================================================

# 1. Background (3.23->6.12)
tf=shapes['TextBox 26'].text_frame; tf.word_wrap=True
bl(tf,'86% of S&P 500 reports contain at least one misleading chart')
bl(tf,'Non-GAAP: #1 SEC concern for 9 years; detection is entirely manual')
bl(tf,'No prior work applies VLMs to automate this detection')

# 2. Motivation (6.92->14.07)
tf_rq=tb(0.40,7.00,6.40,2.10)
hd(tf_rq,'Research Questions')
bl(tf_rq,'RQ1  VLM accuracy across 12 misleader types')
bl(tf_rq,'RQ2  Does textual grounding improve detection?')
bl(tf_rq,'RQ3  Can post-processing tools improve per-type quality?')

try: slide.shapes.add_picture(f'{IMG}/000057_bce3633c432c5bc4_truncated axis.png',
        Inches(0.40),Inches(9.40),Inches(2.80),Inches(1.80))
except: pass
tf_ex=tb(3.30,9.40,3.50,1.80)
wr(tf_ex,'Motivating Example',sz=SZ_BODY,bold=True,color=RED)
wr(tf_ex,'Truncated y-axis (starts at 0.38) exaggerates a 50% gap to appear 5\u00d7 larger.',sz=SZ_CAP,color=DARK)

tf_gap=tb(0.40,11.50,6.40,2.30)
hd(tf_gap,'Research Gap',sz=SZ_BODY)
wr(tf_gap,'No prior work evaluates VLMs for financial chart misleader detection, nor explores tool-augmented pipelines for per-type accuracy.',sz=SZ_BODY,color=BLUE)

# 3. Related Work (14.86->17.83)
tf=shapes['Rectangle 117'].text_frame; tf.word_wrap=True
bl(tf,'Misviz (Tonglet et al., 2025): 2,604 charts, 12 types',sz=SZ_CAP)
bl(tf,'Pandey et al., 2025: VLMs max 30% element detection',sz=SZ_CAP)
bl(tf,'MisVisFix (2025): LLM chart correction, web domain only',sz=SZ_CAP)
wr(tf,'',sz=4)
wr(tf,'References',sz=SZ_CAP,bold=True,color=BLUE)
wr(tf,'[1] Tonglet et al., AAAI 2025  [2] Pandey et al., 2025  [3] MisVisFix, 2025',sz=12,color=GRAY)

# ================================================================
# CENTER COLUMN — METHODOLOGY
# Fig 1: Experimental design (hero, fills column)
# Bottom: Dataset & Models text
# ================================================================

# LaTeX exp design. Measured: ratio=0.5788, h=13.90 -> w=8.05, centered in 9.20
try: slide.shapes.add_picture(f'{F}/exp_design_latex.png',
        Inches(7.67),Inches(3.30),Inches(8.05),Inches(13.90))
except: pass

cap1=tb(7.20,17.25,9.20,0.40)
cap(cap1,'Fig 1. Three-phase experimental design: baseline, pipeline ablation, and SEC validation.')

# ================================================================
# RIGHT COLUMN — RESULTS (single LaTeX-rendered figure)
# ================================================================

# Hide template sub-headers — LaTeX figure has its own
hide('Rectangle 69')
hide('Rectangle 71')
hide('TextBox 288')

# LaTeX results figure. Measured: ratio=0.5212
# At w=7.00: h=7.00/0.5212=13.43
try: slide.shapes.add_picture(f'{F}/results_latex.png',
        Inches(16.80),Inches(3.30),Inches(7.00),Inches(13.43))
except: pass

output='C:/Users/chntw/Documents/7180/DD_v1/FinChartAudit_Poster_v29.pptx'
prs.save(output); print(f'Saved: {output}')
