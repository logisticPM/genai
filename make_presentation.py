"""Create 4+1 slide presentation using presentationV1 template."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation('C:/Users/chntw/Documents/7180/DD_v1/presentationV1 1.pptx')

# Layout indices
LY_TITLE = 4    # TitleSlide1c
LY_CONTENT = 6  # Title and Content_White
LY_TWO_COL = 8  # Two Column_a
LY_IMG = 12     # Content with Image_a
LY_END = 19     # End Slide_Light

BLUE = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
PURPLE = RGBColor(0x75, 0x58, 0x82)

fig = 'C:/Users/chntw/Documents/7180/DD_v1/results/figures'
sec_img = 'C:/Users/chntw/Documents/7180/DD_v1/data/sec_charts/STZ'

# Remove existing slides (delete in reverse)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

def add_slide(layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def set_title(slide, text, size=24):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = BLUE
            return shape

def get_content(slide):
    """Get the content placeholder (idx != 0)."""
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx != 0 and idx != 8:  # not title, not slide number
            return shape
    return None

def add_bullet(tf, text, size=14, color=DARK, bold=False, level=0):
    if tf.paragraphs and tf.paragraphs[0].text == '' and len(tf.paragraphs) == 1:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.level = level
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return p

def add_spacer(tf, size=6):
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = ''
    r.font.size = Pt(size)


# ════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════
s1 = add_slide(LY_TITLE)
set_title(s1, 'FinChartAudit: Detecting Misleading\nFinancial Charts with Vision-Language Models', size=28)
sub = get_content(s1)
if sub:
    sub.text_frame.clear()
    add_bullet(sub.text_frame, 'Mengshan Li  |  Wenshuang Zhou  |  Tong Wu', size=16, bold=True)
    add_bullet(sub.text_frame, 'Khoury College of Computer Sciences, Northeastern University', size=12, color=GRAY)


# ════════════════════════════════════════════════════
# SLIDE 2: Background & RQs
# ════════════════════════════════════════════════════
s2 = add_slide(LY_CONTENT)
set_title(s2, 'Background & Research Questions')
content = get_content(s2)
if content:
    tf = content.text_frame
    tf.clear()

    add_bullet(tf, 'Problem', size=18, bold=True, color=BLUE)
    add_bullet(tf, '86% of S&P 500 reports contain misleading charts (truncated axes, Non-GAAP without GAAP)', size=14, level=1)
    add_bullet(tf, 'Non-GAAP: #1 SEC concern for 9 years; detection entirely manual', size=14, level=1)
    add_bullet(tf, 'No prior work applies VLMs to financial chart misleader detection', size=14, level=1)

    add_spacer(tf)
    add_bullet(tf, 'Research Questions', size=18, bold=True, color=BLUE)
    add_bullet(tf, 'RQ1  VLM accuracy across 12 misleader types (Claude vs Qwen)', size=14, level=1)
    add_bullet(tf, 'RQ2  Does textual grounding (bbox) improve detection?', size=14, level=1)
    add_bullet(tf, 'RQ3  Can post-processing tools improve per-type quality?', size=14, level=1)

    add_spacer(tf)
    add_bullet(tf, 'Key Observation', size=18, bold=True, color=RED)
    add_bullet(tf, 'Binary F1 = 83% masks poor per-type F1 = 39.3% (42% FP on clean charts)', size=14, color=RED, level=1)
    add_bullet(tf, '\u2192 Motivates tool-augmented post-processing (RQ3)', size=14, color=BLUE, level=1)


# ════════════════════════════════════════════════════
# SLIDE 3: Methodology (with LaTeX diagram)
# ════════════════════════════════════════════════════
s3 = add_slide(LY_CONTENT)
set_title(s3, 'Methodology')
# Clear content placeholder text
content = get_content(s3)
if content:
    content.text_frame.clear()
    add_bullet(content.text_frame, '', size=1)  # minimal placeholder

# Add experimental design diagram as image (fills most of slide)
try:
    s3.shapes.add_picture(
        f'{fig}/exp_design_latex.png',
        Inches(1.5), Inches(1.6), Inches(5.0), Inches(5.5))
except: pass

# Add key points on the right
from pptx.util import Emu
tb = s3.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.2))
tb.text_frame.word_wrap = True
tf = tb.text_frame

add_bullet(tf, 'Phase 1: Baseline', size=16, bold=True, color=BLUE)
add_bullet(tf, '2\u00d72 factorial on 271 Misviz charts', size=13, level=1)
add_bullet(tf, 'Claude Haiku-4.5 / Qwen3-VL-8B', size=13, level=1)
add_bullet(tf, 'vision-only / vision + bbox text', size=13, level=1)

add_spacer(tf)
add_bullet(tf, 'Phase 2: Pipeline Ablation', size=16, bold=True, color=GREEN)
add_bullet(tf, '6 strategies: injection vs post-processing', size=13, level=1)
add_bullet(tf, 'Tools: RapidOCR, DePlot, ViT-B/16', size=13, level=1)

add_spacer(tf)
add_bullet(tf, 'Phase 3: SEC Validation', size=16, bold=True, color=PURPLE)
add_bullet(tf, 'Best pipeline on 9 STZ 10-K charts', size=13, level=1)
add_bullet(tf, 'vs. SEC comment letter ground truth', size=13, level=1)


# ════════════════════════════════════════════════════
# SLIDE 4: Results
# ════════════════════════════════════════════════════
s4 = add_slide(LY_CONTENT)
set_title(s4, 'Results')
content = get_content(s4)
if content:
    tf = content.text_frame
    tf.clear()

    add_bullet(tf, 'RQ1/2: Baseline', size=16, bold=True, color=BLUE)
    add_bullet(tf, 'Claude 83.0% vs Qwen 72.3% (+10.7 pp)', size=13, color=GREEN, level=1)
    add_bullet(tf, 'Vision+text: precision \u2191, recall \u2193, marginal net effect', size=13, level=1)

    add_spacer(tf, 4)
    add_bullet(tf, 'RQ3: Pipeline Ablation (Per-Type F1)', size=16, bold=True, color=GREEN)
    add_bullet(tf, 'Tool injection degrades: OCR \u22120.4 pp, DePlot \u22121.7 pp', size=13, color=RED, level=1)
    add_bullet(tf, 'Post-processing improves: CLEAN Veto +4.2 pp', size=13, color=GREEN, level=1)
    add_bullet(tf, 'Best: CLEAN Veto + Classifier = 45.2% (+5.9 pp)', size=13, color=GREEN, bold=True, level=1)
    add_bullet(tf, 'Self-consistency fails: 3-vote increases FP (systematic errors)', size=13, level=1)

    add_spacer(tf, 4)
    add_bullet(tf, 'SEC Case Study: STZ (Constellation Brands)', size=16, bold=True, color=PURPLE)
    add_bullet(tf, 'Baseline: Precision=0.33, F1=0.50 (6 FP on 9 charts)', size=13, level=1)
    add_bullet(tf, 'Pipeline: Precision=1.00, F1=0.80 (0 FP) \u2014 all FP eliminated', size=13, color=GREEN, bold=True, level=1)

# Add results table
try:
    s4.shapes.add_picture(
        f'{fig}/results_table_latex.png',
        Inches(7.0), Inches(1.8), Inches(5.8), Inches(1.5))
except: pass

# Add SEC chart examples
try:
    s4.shapes.add_picture(
        f'{sec_img}/stz-20240229_g22.jpg',
        Inches(7.0), Inches(3.5), Inches(2.8), Inches(1.4))
except: pass
try:
    s4.shapes.add_picture(
        f'{sec_img}/stz-20240229_g10.jpg',
        Inches(10.0), Inches(3.5), Inches(2.8), Inches(1.4))
except: pass

# Labels for SEC images
tb2 = s4.shapes.add_textbox(Inches(7.0), Inches(4.9), Inches(5.8), Inches(0.3))
tb2.text_frame.word_wrap = True
p = tb2.text_frame.paragraphs[0]
r = p.add_run()
r.text = 'Truncated axis: Y from $11,500          Truncated axis: Y range 42\u201345%'
r.font.size = Pt(10)
r.font.color.rgb = GRAY
r.font.italic = True


# ════════════════════════════════════════════════════
# SLIDE 5: Conclusion + Thank You
# ════════════════════════════════════════════════════
s5 = add_slide(LY_CONTENT)
set_title(s5, 'Conclusion & Future Work')
content = get_content(s5)
if content:
    tf = content.text_frame
    tf.clear()

    add_bullet(tf, 'Key Findings', size=18, bold=True, color=BLUE)
    add_bullet(tf, 'Post-processing verifiers improve PT F1 by +5.9 pp (+15%)', size=14, color=GREEN, level=1)
    add_bullet(tf, 'Tool injection into prompt degrades performance by \u22121.7 pp', size=14, color=RED, level=1)
    add_bullet(tf, 'Claude outperforms Qwen by +10.7 pp; text context: marginal net effect', size=14, level=1)
    add_bullet(tf, 'VLM errors are systematic (attention dilution) \u2014 self-consistency fails', size=14, level=1)
    add_bullet(tf, 'SEC: pipeline eliminates all FP while detecting truncated axes', size=14, level=1)

    add_spacer(tf)
    add_bullet(tf, 'Future Work', size=18, bold=True, color=BLUE)
    add_bullet(tf, 'Fine-tune VLM on financial chart domain', size=14, level=1)
    add_bullet(tf, 'Domain adaptation for synth\u2192real classifier gap', size=14, level=1)
    add_bullet(tf, 'Scale SEC coverage with chart-specific dataset', size=14, level=1)

    add_spacer(tf)
    add_spacer(tf)
    add_bullet(tf, 'Thank you! Questions?', size=20, bold=True, color=BLUE)


output = 'C:/Users/chntw/Documents/7180/DD_v1/FinChartAudit_Presentation_v2.pptx'
prs.save(output)
print(f'Saved: {output}')
